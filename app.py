import streamlit as st
import os
import re
import zipfile
import fitz  # PyMuPDF
import pytesseract
import camelot
import pdfkit
import pandas as pd
from pptx import Presentation
from pdf2image import convert_from_path
from PIL import Image
from bs4 import BeautifulSoup
from ebooklib import epub, ITEM_DOCUMENT
import pypandoc

# ==============================================================================
# PART 1: ALL HELPER & PROCESSING FUNCTIONS
# ==============================================================================
# These are the same powerful functions we built in the Colab notebook.

# --- UTILITY: Creates safe filenames ---
def create_safe_filename(text, max_length=50):
    text = text[:max_length]
    text = re.sub(r'[\\/*?:"<>|]', "", text)
    text = re.sub(r'\s+', '_', text)
    text = re.sub(r'_+', '_', text)
    text = text.strip('_ ')
    if not text:
        text = "Untitled"
    return text

# --- FEATURE: OCR to a single Searchable PDF ---
def ocr_to_searchable_pdf(pdf_path, dpi=200):
    base_name = os.path.splitext(pdf_path)[0]
    output_pdf = f"{base_name}_searchable.pdf"
    st.write(f"🔍 Running OCR...")
    
    pages = convert_from_path(pdf_path, dpi=dpi)
    pdf_pages = [pytesseract.image_to_pdf_or_hocr(page, extension='pdf') for page in pages]
    
    with st.spinner(f"Processing {len(pages)} pages..."):
        with open(output_pdf, "w+b") as f: f.write(b"".join(pdf_pages))
    
    st.success(f"✅ Searchable PDF created: {os.path.basename(output_pdf)}")
    return output_pdf

# --- FEATURE: Extract Tables to CSV ---
def extract_tables_from_pdf(pdf_path):
    st.write(f"📊 Extracting tables...")
    try:
        tables = camelot.read_pdf(pdf_path, pages='all', flavor='lattice')
        if not tables.n:
            st.warning("  > No tables detected.")
            return []
        csv_files = []
        for i, table in enumerate(tables):
            csv_file = f"{os.path.splitext(pdf_path)[0]}_table_{i+1}.csv"
            table.to_csv(csv_file)
            csv_files.append(csv_file)
            st.success(f"✅ Saved table {i+1} → {os.path.basename(csv_file)}")
        return csv_files
    except Exception as e:
        st.error(f"❌ Table extraction failed: {e}")
        return []

# --- FEATURE: Extract Images ---
def extract_images_from_pdf(pdf_path):
    st.write(f"🖼️ Extracting images...")
    doc, img_files = fitz.open(pdf_path), []
    for page_num in range(len(doc)):
        for img_index, img in enumerate(doc.get_page_images(page_num)):
            xref = img[0]
            base_image = doc.extract_image(xref)
            img_bytes, img_ext = base_image["image"], base_image["ext"]
            img_filename = f"{os.path.splitext(pdf_path)[0]}_page{page_num+1}_img{img_index+1}.{img_ext}"
            with open(img_filename, "wb") as img_file: img_file.write(img_bytes)
            img_files.append(img_filename)
    doc.close()
    if not img_files: st.warning("  > No images found.")
    else:
        for f in img_files: st.success(f"✅ Extracted image → {os.path.basename(f)}")
    return img_files

# --- FEATURE: Smart Split by Headings ---
def smart_split_pdf_by_headings(pdf_path, headings, chapters_per_split=1):
    st.write(f"📖 Smart splitting into chunks of {chapters_per_split} chapters...")
    doc, split_files = fitz.open(pdf_path), []
    num_splits = (len(headings) + chapters_per_split - 1) // chapters_per_split
    for i in range(num_splits):
        start_idx = i * chapters_per_split
        end_idx = min((i + 1) * chapters_per_split, len(headings))
        start_page = headings[start_idx]['page']
        end_page = headings[end_idx]['page'] - 1 if end_idx < len(headings) else (len(doc) - 1)
        new_doc = fitz.open()
        new_doc.insert_pdf(doc, from_page=start_page, to_page=end_page)
        first_chapter_title = headings[start_idx]['title']
        safe_title = create_safe_filename(first_chapter_title)
        output_pdf = f"Part_{i+1}_{safe_title}.pdf"
        if new_doc.page_count > 0:
            new_doc.save(output_pdf); split_files.append(output_pdf)
            st.success(f"✅ Saved split → {os.path.basename(output_pdf)}")
        new_doc.close()
    doc.close()
    return split_files

# --- FEATURE: Clean Headers and Footers ---
def clean_headers_and_footers(pdf_path):
    st.write(f"🧹 Intelligently cleaning headers and footers...")
    doc = fitz.open(pdf_path)
    full_clean_text = ""
    header_margin, footer_margin = 0.15, 0.85
    noise_patterns = [re.compile(r'page\s*\d+', re.IGNORECASE), re.compile(r'^\s*\d+\s*$', re.IGNORECASE)]
    valuable_patterns = [re.compile(r'https?://\S+'), re.compile(r'\[\d+\]')]
    with st.spinner(f"Cleaning {len(doc)} pages..."):
        for page_num, page in enumerate(doc, 1):
            page_text, page_height = "", page.rect.height
            blocks = page.get_text("dict")["blocks"]
            for block in blocks:
                if "lines" in block:
                    for line in block["lines"]:
                        line_text = " ".join([span["text"] for span in line["spans"]]).strip()
                        y_position = line["bbox"][1]
                        is_header, is_footer = y_position < (page_height * header_margin), y_position > (page_height * footer_margin)
                        is_noise = any(p.search(line_text) for p in noise_patterns)
                        is_valuable = any(p.search(line_text) for p in valuable_patterns)
                        if is_valuable or not (is_header or is_footer) or ((is_header or is_footer) and not is_noise):
                            page_text += line_text + "\n"
            full_clean_text += f"\n\n--- Page {page_num} ---\n\n{page_text}"
    doc.close()
    output_txt = f"{os.path.splitext(pdf_path)[0]}_cleaned.txt"
    with open(output_txt, "w", encoding="utf-8") as f: f.write(full_clean_text)
    st.success(f"✅ Cleaning complete → {os.path.basename(output_txt)}")
    return output_txt

# --- FEATURE: Convert PowerPoint to Text ---
def convert_pptx_to_text(pptx_path):
    st.write(f"📝 Extracting text from presentation...")
    try:
        prs = Presentation(pptx_path)
        full_text = "\n".join([shape.text for slide in prs.slides for shape in slide.shapes if hasattr(shape, "text")])
        output_txt = f"{os.path.splitext(pptx_path)[0]}.txt"
        with open(output_txt, "w", encoding="utf-8") as f: f.write(full_text)
        return output_txt
    except Exception as e:
        st.error(f"❌ PowerPoint extraction failed: {e}")
        return None

# --- FEATURE: Convert Spreadsheet to PDF ---
def convert_spreadsheet_to_pdf(file_path):
    st.write(f"📊 Converting spreadsheet to PDF...")
    base_name, ext = os.path.splitext(file_path)
    output_pdf = f"{base_name}.pdf"
    try:
        df = pd.read_excel(file_path) if ext.lower() in ['.xlsx', '.xls'] else pd.read_csv(file_path)
        html_table = df.to_html(index=False, border=1)
        pdfkit.from_string(html_table, output_pdf)
        return output_pdf
    except Exception as e:
        st.error(f"❌ Spreadsheet conversion failed: {e}")
        return None

# ==============================================================================
# PART 2: STREAMLIT USER INTERFACE
# ==============================================================================

st.set_page_config(layout="wide")
st.title(" универсальный препроцессор контента 🤖")
st.write("Upload any supported file to convert or process it for AI analysis.")

uploaded_file = st.file_uploader("Choose a file to begin...")

if uploaded_file is not None:
    # Save the uploaded file to a temporary location
    with open(uploaded_file.name, "wb") as f:
        f.write(uploaded_file.getbuffer())
    
    filename = uploaded_file.name
    base_name, ext = os.path.splitext(filename)

    # --- WORKFLOW 1: Simple Files (DOCX, MD, TXT) ---
    if ext.lower() in ['.docx', '.md', '.txt']:
        st.subheader("Convert Document")
        chosen_format = st.radio("Convert to:", ('PDF', 'TXT'), horizontal=True)
        
        if st.button("Convert and Download"):
            with st.spinner(f"Converting {filename} to {chosen_format}..."):
                pandoc_format = 'plain' if chosen_format == 'TXT' else 'pdf'
                output_filename = f"{base_name}.{chosen_format.lower()}"
                pypandoc.convert_file(filename, pandoc_format, outputfile=output_filename, extra_args=['--pdf-engine=wkhtmltopdf'])
                
                with open(output_filename, "rb") as f:
                    st.download_button(f"Download {os.path.basename(output_filename)}", f, file_name=os.path.basename(output_filename))
                os.remove(output_filename) # Clean up

    # --- WORKFLOW 2: PowerPoint & Spreadsheets ---
    elif ext.lower() in ['.pptx', '.ppt']:
        st.subheader("Process PowerPoint")
        if st.button("Extract Text and Download"):
            with st.spinner("Extracting text..."):
                output_file = convert_pptx_to_text(filename)
                if output_file:
                    with open(output_file, "rb") as f:
                        st.download_button("Download .txt", f, file_name=os.path.basename(output_file))
                    os.remove(output_file)

    elif ext.lower() in ['.xlsx', '.xls', '.csv']:
        st.subheader("Process Spreadsheet")
        if st.button("Convert to PDF and Download"):
            with st.spinner("Converting to PDF..."):
                output_file = convert_spreadsheet_to_pdf(filename)
                if output_file:
                    with open(output_file, "rb") as f:
                        st.download_button("Download .pdf", f, file_name=os.path.basename(output_file))
                    os.remove(output_file)

    # --- WORKFLOW 3: Advanced PDF Options ---
    elif ext.lower() == '.pdf':
        st.subheader("Advanced PDF Processing")
        
        doc = fitz.open(filename)
        headings, last_heading_title = [], ""
        heading_pattern = r'^(Chapter\s+\d+.*|Section\s+\d+.*)'
        for page in doc:
            text = page.get_text("text")
            for line in text.split('\n'):
                match = re.search(heading_pattern, line, flags=re.IGNORECASE)
                current_heading_title = match.group(0).strip() if match else ""
                if match and current_heading_title != last_heading_title:
                    headings.append({'title': current_heading_title, 'page': page.number})
                    last_heading_title = current_heading_title
                    break 
        doc.close()
        
        col1, col2 = st.columns(2)
        with col1:
            st.write("Extraction Options:")
            ocr_check = st.checkbox("Create Searchable PDF (OCR)", help="For scanned, image-only PDFs.")
            table_check = st.checkbox("Extract Tables into CSVs", value=True)
            image_check = st.checkbox("Extract Images", value=True)
            clean_check = st.checkbox("Clean Headers & Footers", value=False)
        
        with col2:
            st.write("Splitting Options:")
            smart_split_check = st.checkbox("Smart Split by Chapter/Section", value=True)
            if headings:
                selected_split_value = st.slider("Group Chapters/Sections by:", min_value=1, max_value=len(headings), value=1)
            else:
                st.info("No 'Chapter X' or 'Section X' headings found for smart splitting.")
                selected_split_value = 1

        if st.button("Process PDF"):
            with st.spinner("Processing PDF with selected options..."):
                all_files = []
                if ocr_check: all_files.append(ocr_to_searchable_pdf(filename))
                if table_check: all_files.extend(extract_tables_from_pdf(filename))
                if image_check: all_files.extend(extract_images_from_pdf(filename))
                if smart_split_check:
                    if not headings:
                        st.warning("Smart Split selected, but no headings were found.")
                    else:
                        all_files.extend(smart_split_pdf_by_headings(filename, headings, selected_split_value))
                if clean_check: all_files.append(clean_headers_and_footers(filename))
                
                if not all_files:
                    st.warning("No actions were selected or no files were generated.")
                else:
                    zip_filename = f"{base_name}_processed.zip"
                    with zipfile.ZipFile(zip_filename, 'w') as zipf:
                        for f in all_files:
                            if os.path.exists(f): zipf.write(f); os.remove(f)
                    
                    with open(zip_filename, "rb") as f:
                        st.download_button(f"Download Results ({os.path.basename(zip_filename)})", f, file_name=zip_filename)
                    os.remove(zip_filename)

    # --- WORKFLOW 4: EPUB Splitting ---
    elif ext.lower() == '.epub':
        st.subheader("Split EPUB")
        book = epub.read_epub(filename)
        items = list(book.get_items_of_type(ITEM_DOCUMENT))
        chapters = [BeautifulSoup(item.get_content(), 'html.parser').get_text().strip() for item in items if item]
        total_chapters = len(chapters)
        
        selected_split_value = st.slider("Group Chapters by:", min_value=1, max_value=total_chapters, value=1)
        
        if st.button("Split EPUB to PDFs"):
            with st.spinner(f"Splitting EPUB into chunks of {selected_split_value} chapters..."):
                split_files = []
                total_splits = (total_chapters + selected_split_value - 1) // selected_split_value
                for i in range(total_splits):
                    start_idx, end_idx = i * selected_split_value, min((i + 1) * selected_split_value, total_chapters)
                    split_chapters = chapters[start_idx:end_idx]
                    first_chapter_content = split_chapters[0]
                    chapter_title = "Untitled_Chapter"
                    for line in first_chapter_content.split('\n')[:10]:
                        line = line.strip()
                        if 0 < len(line) < 70: chapter_title = line; break
                    safe_title = create_safe_filename(chapter_title)
                    output_pdf = f"Part_{i+1}_{safe_title}.pdf"
                    html_content = "".join([f"<h1>Chapter {j+1}</h1><p>{chap}</p>" for j, chap in enumerate(split_chapters, start=start_idx)])
                    pdfkit.from_string(html_content, output_pdf, options={"enable-local-file-access": ""})
                    split_files.append(output_pdf)
                
                zip_filename = f"{base_name}_split.zip"
                with zipfile.ZipFile(zip_filename, 'w') as zipf:
                    for f in split_files: zipf.write(f); os.remove(f)
                
                with open(zip_filename, "rb") as f:
                    st.download_button(f"Download Results ({os.path.basename(zip_filename)})", f, file_name=zip_filename)
                os.remove(zip_filename)

    # Clean up the original uploaded file
    if os.path.exists(filename):
        os.remove(filename)